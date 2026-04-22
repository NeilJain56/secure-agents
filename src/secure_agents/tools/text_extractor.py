"""Text extraction tool for the document sorting pipeline.

Extracts text from PDF, DOCX, DOC, PPTX, and XLSX files.  Unlike the
``document_parser`` tool (which runs inside Docker for untrusted email
attachments), this tool runs directly on the host because the files are
the user's own trusted local documents.  File-type and size validation
is still applied as defense in depth.
"""

from __future__ import annotations

import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any

import structlog

from secure_agents.core.base_tool import BaseTool
from secure_agents.core.registry import register_tool

logger = structlog.get_logger()

_SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".doc", ".pptx", ".xlsx"}

# Extraction limits — enough text for Jaccard pre-filtering AND LLM calls,
# without spending minutes reading a 200-page PDF or a 20 MB spreadsheet.
_MAX_PDF_PAGES = 15       # first N pages of a PDF
_MAX_XLSX_ROWS = 300      # first N rows per sheet in an XLSX

# Magic byte prefixes for quick sanity checks
_MAGIC: dict[str, list[bytes]] = {
    ".pdf":  [b"%PDF"],
    ".docx": [b"PK\x03\x04", b"PK\x05\x06"],
    ".pptx": [b"PK\x03\x04", b"PK\x05\x06"],
    ".xlsx": [b"PK\x03\x04", b"PK\x05\x06"],
    ".doc":  [b"\xd0\xcf\x11\xe0"],
}


def _check_magic(path: Path) -> bool:
    """Return True if the file's leading bytes match the expected magic."""
    suffix = path.suffix.lower()
    sigs = _MAGIC.get(suffix)
    if sigs is None:
        return True  # no signature to check
    try:
        header = path.read_bytes()[:16]
    except OSError:
        return False
    return any(header.startswith(s) for s in sigs)


# ── Per-format extractors ────────────────────────────────────────────────────

def _extract_pdf(path: Path) -> str:
    import pdfplumber

    parts: list[str] = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages[:_MAX_PDF_PAGES]:
            text = page.extract_text()
            if text:
                parts.append(text)
    return "\n\n".join(parts)


def _extract_docx(path: Path) -> str:
    import docx

    doc = docx.Document(str(path))
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_doc(path: Path) -> str:
    """Extract text from legacy .doc (OLE2) via textutil (macOS) or
    libreoffice (Linux).  Falls back to a clear error when neither is
    available.
    """
    # macOS ships textutil
    if shutil.which("textutil"):
        with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as tmp:
            tmp_path = Path(tmp.name)
        try:
            subprocess.run(
                ["textutil", "-convert", "txt", "-output", str(tmp_path), str(path)],
                check=True,
                capture_output=True,
                timeout=60,
            )
            return tmp_path.read_text(encoding="utf-8", errors="replace")
        finally:
            tmp_path.unlink(missing_ok=True)

    # Linux fallback: libreoffice headless → txt
    if shutil.which("libreoffice"):
        with tempfile.TemporaryDirectory() as tmpdir:
            subprocess.run(
                [
                    "libreoffice", "--headless", "--convert-to", "txt:Text",
                    "--outdir", tmpdir, str(path),
                ],
                check=True,
                capture_output=True,
                timeout=120,
            )
            txt = Path(tmpdir) / (path.stem + ".txt")
            if txt.exists():
                return txt.read_text(encoding="utf-8", errors="replace")
            raise RuntimeError(f"libreoffice produced no output for {path.name}")

    raise RuntimeError(
        f"Cannot extract text from .doc files: neither 'textutil' (macOS) "
        f"nor 'libreoffice' (Linux) is installed.  Install one of them or "
        f"convert the file to .docx."
    )


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            parts.append(f"[Slide {slide_idx}]\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True, max_row=_MAX_XLSX_ROWS):
            cells = [str(c) if c is not None else "" for c in row]
            line = " | ".join(cells).strip(" |")
            if line:
                rows.append(line)
        if rows:
            parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(parts)


_EXTRACTORS = {
    ".pdf":  _extract_pdf,
    ".docx": _extract_docx,
    ".doc":  _extract_doc,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx,
}


# ── Tool registration ────────────────────────────────────────────────────────

@register_tool("text_extractor")
class TextExtractorTool(BaseTool):
    """Extracts text from PDF, DOCX, DOC, PPTX, and XLSX files.

    Designed for the document-sorting pipeline where files are the user's
    own trusted local documents.  Validates file type and size as defense
    in depth but does **not** route through the Docker sandbox.
    """

    name = "text_extractor"
    description = "Extract text content from PDF, DOCX, DOC, PPTX, and XLSX files"

    def __init__(self, config: dict | None = None) -> None:
        super().__init__(config)
        self.max_file_size_mb: int = self.config.get("max_file_size_mb", 100)

    def execute(self, **kwargs: Any) -> dict:
        """Extract text from a single file.

        kwargs:
            file_path: str — absolute or relative path to the file (required)

        Returns:
            {"text": str, "file_type": str, "filename": str, "size_bytes": int}
            or {"error": str} on failure.
        """
        raw_path = kwargs.get("file_path", "")
        if not raw_path:
            return {"error": "No file_path provided"}

        path = Path(raw_path).resolve()
        if not path.exists() or not path.is_file():
            return {"error": f"File not found: {path.name}"}

        suffix = path.suffix.lower()
        if suffix not in _SUPPORTED_EXTENSIONS:
            return {"error": f"Unsupported file type: {suffix}. Supported: {sorted(_SUPPORTED_EXTENSIONS)}"}

        size_mb = path.stat().st_size / (1024 * 1024)
        if size_mb > self.max_file_size_mb:
            return {"error": f"File too large: {size_mb:.1f}MB (max {self.max_file_size_mb}MB)"}

        if not _check_magic(path):
            return {"error": f"File content does not match {suffix} format (magic byte mismatch)"}

        try:
            extractor = _EXTRACTORS[suffix]
            text = extractor(path)
            logger.info(
                "text_extractor.ok",
                filename=path.name,
                file_type=suffix.lstrip("."),
                chars=len(text),
            )
            return {
                "text": text,
                "file_type": suffix.lstrip("."),
                "filename": path.name,
                "size_bytes": path.stat().st_size,
            }
        except Exception as e:
            logger.error("text_extractor.error", filename=path.name, error=str(e))
            return {"error": f"Extraction failed for {path.name}: {e}"}

    def validate_config(self) -> bool:
        return True
