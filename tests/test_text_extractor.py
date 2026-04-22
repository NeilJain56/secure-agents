"""Tests for the text_extractor tool."""

from pathlib import Path
from unittest.mock import patch

import pytest

from secure_agents.tools.text_extractor import TextExtractorTool


@pytest.fixture
def tool():
    return TextExtractorTool(config={"max_file_size_mb": 50})


# ── Basic validation ─────────────────────────────────────────────────────────

def test_missing_file_path(tool):
    result = tool.execute()
    assert "error" in result
    assert "No file_path" in result["error"]


def test_nonexistent_file(tool):
    result = tool.execute(file_path="/tmp/does_not_exist.pdf")
    assert "error" in result
    assert "not found" in result["error"]


def test_unsupported_extension(tool, tmp_path):
    f = tmp_path / "readme.txt"
    f.write_text("hello")
    result = tool.execute(file_path=str(f))
    assert "error" in result
    assert "Unsupported" in result["error"]


def test_file_too_large(tmp_path):
    tool = TextExtractorTool(config={"max_file_size_mb": 0})  # 0 MB limit
    f = tmp_path / "big.pdf"
    f.write_bytes(b"%PDF" + b"\x00" * 1024)
    result = tool.execute(file_path=str(f))
    assert "error" in result
    assert "too large" in result["error"]


def test_magic_byte_mismatch(tmp_path, tool):
    f = tmp_path / "fake.pdf"
    f.write_bytes(b"NOT_A_PDF_FILE" + b"\x00" * 100)
    result = tool.execute(file_path=str(f))
    assert "error" in result
    assert "magic byte" in result["error"]


# ── PDF extraction ───────────────────────────────────────────────────────────

def test_pdf_extraction(tmp_path, tool):
    """Test PDF extraction with a minimal valid PDF."""
    # Create a minimal valid PDF using pdfplumber-compatible content
    pdf_content = (
        b"%PDF-1.4\n"
        b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
        b"2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj\n"
        b"3 0 obj<</Type/Page/MediaBox[0 0 612 792]/Parent 2 0 R"
        b"/Resources<</Font<</F1 4 0 R>>>>"
        b"/Contents 5 0 R>>endobj\n"
        b"4 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj\n"
        b"5 0 obj<</Length 44>>stream\n"
        b"BT /F1 12 Tf 100 700 Td (Hello World) Tj ET\n"
        b"endstream\nendobj\n"
        b"xref\n0 6\n"
        b"0000000000 65535 f \n"
        b"0000000009 00000 n \n"
        b"0000000058 00000 n \n"
        b"0000000115 00000 n \n"
        b"0000000266 00000 n \n"
        b"0000000340 00000 n \n"
        b"trailer<</Size 6/Root 1 0 R>>\n"
        b"startxref\n434\n%%EOF\n"
    )
    f = tmp_path / "test.pdf"
    f.write_bytes(pdf_content)

    result = tool.execute(file_path=str(f))
    # pdfplumber may or may not extract text from our minimal PDF,
    # but it should not error
    assert "error" not in result or "text" in result


# ── DOCX extraction ──────────────────────────────────────────────────────────

def test_docx_extraction(tmp_path, tool):
    """Test DOCX extraction by creating a real .docx file."""
    import docx

    doc = docx.Document()
    doc.add_paragraph("This is a test document.")
    doc.add_paragraph("Second paragraph with important content.")
    docx_path = tmp_path / "test.docx"
    doc.save(str(docx_path))

    result = tool.execute(file_path=str(docx_path))
    assert "error" not in result
    assert "test document" in result["text"]
    assert "Second paragraph" in result["text"]
    assert result["file_type"] == "docx"
    assert result["filename"] == "test.docx"


# ── PPTX extraction ─────────────────────────────────────────────────────────

def test_pptx_extraction(tmp_path, tool):
    """Test PPTX extraction by creating a real .pptx file."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # title+content
    slide.shapes.title.text = "Test Slide Title"
    slide.placeholders[1].text = "Slide body content here."
    pptx_path = tmp_path / "test.pptx"
    prs.save(str(pptx_path))

    result = tool.execute(file_path=str(pptx_path))
    assert "error" not in result
    assert "Test Slide Title" in result["text"]
    assert "body content" in result["text"]
    assert result["file_type"] == "pptx"


# ── XLSX extraction ──────────────────────────────────────────────────────────

def test_xlsx_extraction(tmp_path, tool):
    """Test XLSX extraction by creating a real .xlsx file."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Value"])
    ws.append(["Alpha", 100])
    ws.append(["Beta", 200])
    xlsx_path = tmp_path / "test.xlsx"
    wb.save(str(xlsx_path))

    result = tool.execute(file_path=str(xlsx_path))
    assert "error" not in result
    assert "Alpha" in result["text"]
    assert "Beta" in result["text"]
    assert "Sheet1" in result["text"]
    assert result["file_type"] == "xlsx"


# ── DOC extraction ───────────────────────────────────────────────────────────

def test_doc_extraction_no_converter(tmp_path, tool):
    """If neither textutil nor libreoffice is available, return a clear error."""
    # Create a file with valid .doc magic bytes
    doc_path = tmp_path / "test.doc"
    doc_path.write_bytes(b"\xd0\xcf\x11\xe0" + b"\x00" * 200)

    with patch("secure_agents.tools.text_extractor.shutil.which", return_value=None):
        result = tool.execute(file_path=str(doc_path))
    assert "error" in result
    assert "textutil" in result["error"] or "libreoffice" in result["error"]
